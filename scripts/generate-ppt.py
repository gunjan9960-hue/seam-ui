"""
Seam Capstone Deck Generator
Run: python3 scripts/generate-ppt.py
Output: docs/Seam_Capstone_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x11, 0x17)   # #0F1117
BLUE      = RGBColor(0x4F, 0x6B, 0xF5)   # #4F6BF5
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DIM       = RGBColor(0xA0, 0xA8, 0xB8)   # subdued text
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # pass
RED       = RGBColor(0xF8, 0x71, 0x71)   # fail
YELLOW    = RGBColor(0xFC, 0xD3, 0x4D)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
CARD_BG   = RGBColor(0x1A, 0x1D, 0x27)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide)
    return slide

def fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box

def rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def hline(slide, y, color=RGBColor(0x2A, 0x2D, 0x3A)):
    line = slide.shapes.add_shape(1, Inches(0), y, W, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def dot(slide, x, y, r=6):
    circ = slide.shapes.add_shape(9, x - Emu(r*914), y - Emu(r*914),
                                  Emu(r*914*2), Emu(r*914*2))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()

def phase_tag(slide, label, color=BLUE):
    box = slide.shapes.add_shape(1, Inches(0.55), Inches(0.28), Inches(2.4), Inches(0.38))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size  = Pt(10)
    run.font.bold  = True
    run.font.color.rgb = WHITE

def bullet_block(slide, items, x, y, w, size=14, color=DIM, gap=Inches(0.38)):
    for item in items:
        txbox(slide, item, x, y, w, Inches(0.5), size=size, color=color)
        y += gap
    return y

def card(slide, title, body_lines, x, y, w, h,
         title_color=WHITE, accent=BLUE):
    # card bg
    c = slide.shapes.add_shape(1, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = RGBColor(0x2A, 0x2D, 0x3A)
    # accent bar
    bar = slide.shapes.add_shape(1, x, y, Pt(3), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # title
    txbox(slide, title, x + Inches(0.22), y + Inches(0.15), w - Inches(0.3), Inches(0.35),
          size=12, bold=True, color=title_color)
    # body
    body = "\n".join(body_lines)
    txbox(slide, body, x + Inches(0.22), y + Inches(0.52), w - Inches(0.3), h - Inches(0.6),
          size=11, color=DIM)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────
prs = new_prs()

# ── S1: TITLE ─────────────────────────────────────────────────────────────────
s = blank_slide(prs)
# big seam wordmark
txbox(s, "seam.", Inches(1.5), Inches(1.6), Inches(6), Inches(1.4),
      size=80, bold=True, color=WHITE)
dot(s, Inches(4.72), Inches(1.95), r=16)
txbox(s, "AI Search Layer for Product Managers",
      Inches(1.5), Inches(3.1), Inches(8), Inches(0.6),
      size=22, color=DIM)
hline(s, Inches(3.85))
txbox(s, "Capstone Project  ·  Cohort 6",
      Inches(1.5), Inches(3.95), Inches(6), Inches(0.45), size=14, color=DIM)

# ── S2: THE PROBLEM ───────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "THE PROBLEM")
txbox(s, "PMs are drowning in context-switching",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.7),
      size=30, bold=True, color=WHITE)
hline(s, Inches(1.65))

stats = [
    ("25 min",    "lost per context switch\nacross tools"),
    ("66%",       "PM time spent on\nmanual information work"),
    ("3–6 wks",   "to build context\nwhen joining a new product"),
]
cx = Inches(0.7)
for val, label in stats:
    txbox(s, val, cx, Inches(2.1), Inches(3.2), Inches(0.8), size=42, bold=True, color=BLUE)
    txbox(s, label, cx, Inches(2.95), Inches(3.2), Inches(0.7), size=13, color=DIM)
    cx += Inches(4.1)

txbox(s, "The answer lives in Notion. The context is in Jira. The escalation is in Slack. The decision is in a Google Doc nobody updated.",
      Inches(0.7), Inches(4.3), Inches(11.5), Inches(0.8),
      size=16, color=DIM, italic=True)

txbox(s, "PM CHECKLIST  ·  Phase 1: Problem Definition",
      Inches(0.7), Inches(5.3), Inches(12), Inches(0.35),
      size=10, bold=True, color=BLUE)
checks = "✓  Quantify the pain with real numbers   ✓  Identify who feels it most   ✓  Confirm it is worth solving"
txbox(s, checks, Inches(0.7), Inches(5.65), Inches(12), Inches(0.5), size=13, color=GREEN)

# ── S3: USER PERSONA ──────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 1 · DISCOVERY", PURPLE)
txbox(s, "Who we are building for",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=30, bold=True, color=WHITE)
hline(s, Inches(1.55))

card(s, "👤  Primary Persona — The B2B SaaS PM",
     ["Name: Priya / Rajan  ·  Exp: 2–6 years in product",
      "Company: Series A–C SaaS startup, 50–300 employees",
      "Tools: Notion + Jira + Slack + Google Docs (all 4)",
      "Pain: Spends 2–3 hrs/day digging for context before decisions",
      "Goal: Answer stakeholder questions confidently in <1 min",
      "Trigger: Joins new product or cross-functional project"],
     Inches(0.55), Inches(1.75), Inches(5.8), Inches(2.8), accent=PURPLE)

card(s, "🧠  Key Insight from Research",
     ["PMs don't need more tools — they need fewer tabs",
      "The bottleneck is retrieval, not reasoning",
      "'I know the answer is somewhere' = 25 min search",
      "Trust the answer only when it shows the source"],
     Inches(6.7), Inches(1.75), Inches(5.9), Inches(2.8), accent=BLUE)

txbox(s, "PM CHECKLIST  ·  Phase 1: User Research",
      Inches(0.55), Inches(5.05), Inches(12), Inches(0.35),
      size=10, bold=True, color=BLUE)
checks = "✓  Define primary persona clearly   ✓  Separate 'nice to have' from 'critical pain'   ✓  Quote users verbatim in the doc"
txbox(s, checks, Inches(0.55), Inches(5.42), Inches(12), Inches(0.5), size=13, color=GREEN)

# ── S4: MARKET SIZING ─────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 1 · MARKET SIZING", PURPLE)
txbox(s, "The market is large — and underserved",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=30, bold=True, color=WHITE)
hline(s, Inches(1.55))

markets = [
    ("TAM", "$4.2B", "All PM productivity tools\nglobally"),
    ("SAM", "$420M", "B2B SaaS PMs with\nNotion+Jira+Slack"),
    ("SOM", "$12M", "India-based B2B SaaS PMs\nat Series A–C startups"),
]
cx = Inches(0.7)
for label, val, desc in markets:
    txbox(s, label, cx, Inches(2.0), Inches(3.5), Inches(0.45), size=11, bold=True, color=DIM)
    txbox(s, val,   cx, Inches(2.4), Inches(3.5), Inches(0.7), size=38, bold=True, color=BLUE)
    txbox(s, desc,  cx, Inches(3.1), Inches(3.5), Inches(0.6), size=13, color=DIM)
    cx += Inches(4.1)

card(s, "Why India B2B SaaS PMs First",
     ["~18,000 PMs at Series A–C SaaS companies",
      "High tool adoption (all 4 tools) = richest corpus",
      "Low existing AI-in-PM-workflow penetration",
      "English-first, tool-native, high willingness to pay ($15–40/mo)"],
     Inches(0.7), Inches(4.5), Inches(11.8), Inches(1.8), accent=YELLOW)

txbox(s, "PM CHECKLIST  ·  Phase 1: Market Sizing",
      Inches(0.55), Inches(6.55), Inches(12), Inches(0.35),
      size=10, bold=True, color=BLUE)
txbox(s, "✓  TAM/SAM/SOM with defensible assumptions   ✓  Know your beachhead segment deeply   ✓  Name 3 competitors and their gap",
      Inches(0.55), Inches(6.9), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S5: COMPETITIVE LANDSCAPE ────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 1 · COMPETITIVE ANALYSIS", PURPLE)
txbox(s, "Competitors exist — but none solve this exactly",
      Inches(0.55), Inches(0.85), Inches(12), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

competitors = [
    ("Notion AI",      "In-doc AI only\nCan't search across Jira/Slack",    RED),
    ("Guru",           "Knowledge base tool\nRequires manual curation",       YELLOW),
    ("Glean",          "Enterprise search\nNo PM-specific intent layer",      YELLOW),
    ("ChatGPT",        "No source attribution\nHallucinates context",         RED),
    ("seam. ✓",        "Cross-tool search\nPM intents · Cited answers",       GREEN),
]
cx = Inches(0.4)
for name, desc, col in competitors:
    box = s.shapes.add_shape(1, cx, Inches(2.0), Inches(2.3), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = col
    tf = box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = name
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = col
    txbox(s, desc, cx + Inches(0.1), Inches(3.1), Inches(2.1), Inches(1.2), size=11, color=DIM, align=PP_ALIGN.CENTER)
    cx += Inches(2.5)

txbox(s, "Key differentiator: Seam is the only tool that combines cross-tool retrieval + PM-intent classification + cited answers in one query.",
      Inches(0.55), Inches(5.0), Inches(12), Inches(0.7), size=15, color=WHITE, italic=True)

txbox(s, "PM CHECKLIST  ·  Phase 1: Competitive Analysis",
      Inches(0.55), Inches(6.0), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  2x2 positioning matrix   ✓  Name the exact gap you exploit   ✓  Know why leaders haven't filled it",
      Inches(0.55), Inches(6.35), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S6: MVP DEFINITION ───────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 2 · MVP DEFINITION", BLUE)
txbox(s, "What we decided to build — and what we cut",
      Inches(0.55), Inches(0.85), Inches(12), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

card(s, "✅  IN SCOPE (MVP)",
     ["Cross-tool search: Notion, Jira, Slack, Google Docs",
      "7 PM-intent types with dynamic follow-up suggestions",
      "Single cited answer with source cards + freshness warning",
      "Ambiguity detection for short/vague queries",
      "Conversation history (multi-turn follow-ups)"],
     Inches(0.55), Inches(1.75), Inches(6.0), Inches(2.9), accent=GREEN)

card(s, "❌  OUT OF SCOPE (MVP)",
     ["Real-time connector sync (static corpus for Phase 4)",
      "Native mobile app",
      "User auth / multi-workspace support",
      "PRD generator (Phase 3 feature, deprioritised)",
      "Slack bot / Jira plugin"],
     Inches(6.9), Inches(1.75), Inches(5.8), Inches(2.9), accent=RED)

txbox(s, "PM CHECKLIST  ·  Phase 2: Scoping",
      Inches(0.55), Inches(5.1), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  Explicit 'In/Out of scope' list   ✓  Each cut item has a reason   ✓  MVP delivers the core value loop end-to-end",
      Inches(0.55), Inches(5.45), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S7: SUCCESS METRICS ──────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 2 · SUCCESS METRICS", BLUE)
txbox(s, "How we know if we shipped the right thing",
      Inches(0.55), Inches(0.85), Inches(12), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

metrics = [
    ("Answer Relevance",  "≥ 85%",       "Does it answer the question?",         "4.8/5 ✓",  GREEN),
    ("Citation Accuracy", "≥ 90%",       "Are claims backed by sources?",         "5.0/5 ✓",  GREEN),
    ("Hallucination",     "< 5%",        "Are all facts from indexed docs?",      "4.9/5 ✓",  GREEN),
    ("Total Latency",     "< 10 sec",    "Time from query to full answer",        "8.7s ✓",   GREEN),
    ("Gen Latency",       "< 6 sec",     "LLM generation time alone",             "7.3s ✗",   RED),
]
y = Inches(1.8)
txbox(s, "Metric", Inches(0.55), y, Inches(2.8), Inches(0.35), size=10, bold=True, color=DIM)
txbox(s, "Target", Inches(3.4),  y, Inches(1.5), Inches(0.35), size=10, bold=True, color=DIM)
txbox(s, "Definition", Inches(5.0), y, Inches(4.5), Inches(0.35), size=10, bold=True, color=DIM)
txbox(s, "v1 Result", Inches(9.5), y, Inches(2.5), Inches(0.35), size=10, bold=True, color=DIM)
y += Inches(0.4)
hline(s, y)
y += Inches(0.05)

for name, target, defn, result, col in metrics:
    txbox(s, name,   Inches(0.55), y, Inches(2.8), Inches(0.38), size=13, bold=True, color=WHITE)
    txbox(s, target, Inches(3.4),  y, Inches(1.5), Inches(0.38), size=13, color=BLUE)
    txbox(s, defn,   Inches(5.0),  y, Inches(4.5), Inches(0.38), size=12, color=DIM)
    txbox(s, result, Inches(9.5),  y, Inches(2.5), Inches(0.38), size=13, bold=True, color=col)
    y += Inches(0.42)

txbox(s, "PM CHECKLIST  ·  Phase 2: Metrics",
      Inches(0.55), Inches(6.55), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  Define metrics BEFORE building   ✓  Each metric is measurable   ✓  Know who owns each metric",
      Inches(0.55), Inches(6.9), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S8: ARCHITECTURE ─────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 3 · ARCHITECTURE", RGBColor(0x34,0xD3,0x99))
txbox(s, "How Seam works under the hood",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

steps = [
    ("1", "Query In",        "PM types a question\nin natural language",         BLUE),
    ("2", "Intent Detect",   "7 PM intents classified\n(decision, spec, customer…)", PURPLE),
    ("3", "Hybrid RAG",      "BM25 + Voyage AI embeddings\n+ RRF + Re-ranking",  RGBColor(0x60,0xA5,0xFA)),
    ("4", "Generate",        "Claude Sonnet 4.6\nwith intent-aware prompt",      YELLOW),
    ("5", "Answer Out",      "Cited answer + source cards\n+ freshness warning", GREEN),
]
cx = Inches(0.4)
for num, title, desc, col in steps:
    box = s.shapes.add_shape(1, cx, Inches(2.0), Inches(2.3), Inches(1.9))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = col
    txbox(s, num,   cx + Inches(0.15), Inches(2.1), Inches(0.4), Inches(0.4), size=22, bold=True, color=col)
    txbox(s, title, cx + Inches(0.15), Inches(2.55), Inches(2.0), Inches(0.35), size=13, bold=True, color=WHITE)
    txbox(s, desc,  cx + Inches(0.15), Inches(2.92), Inches(2.0), Inches(0.9),  size=11, color=DIM)
    if num != "5":
        txbox(s, "→", cx + Inches(2.35), Inches(2.6), Inches(0.3), Inches(0.4), size=18, color=DIM)
    cx += Inches(2.5)

card(s, "Key Technical Decisions & Why",
     ["Hybrid RAG (BM25 + embeddings) > pure embeddings — better keyword recall for PM terms",
      "Voyage AI rerank-2 — cuts noise, especially for ambiguous 1-2 word queries",
      "Intent-aware prompts — reduces answer length by 30%, improves relevance",
      "43 documents indexed — enough for a real PM workflow (not a toy demo)"],
     Inches(0.55), Inches(4.2), Inches(12.1), Inches(2.0), accent=BLUE)

txbox(s, "PM CHECKLIST  ·  Phase 3: Architecture Decisions",
      Inches(0.55), Inches(6.55), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  Know why each tech choice was made   ✓  Document trade-offs   ✓  Be able to explain it in 60 seconds",
      Inches(0.55), Inches(6.9), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S9: RISK REGISTER ────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 3 · RISK REGISTER", YELLOW)
txbox(s, "Every PM must own the risk register",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

risks = [
    ("A1", "Hallucination",           "Answer contains facts not in sources",  "Verified against full docs in eval",   GREEN),
    ("A2", "Wrong Retrieval",         "Correct intent, wrong docs retrieved",   "Hybrid RAG + re-rank reduces this",   YELLOW),
    ("A3", "Intent Misclassification","'Build vs buy' parsed as wrong intent",  "Add ambiguity detection (v2 fix)",    YELLOW),
    ("A4", "Stale Data",              "Docs not synced, answer is outdated",    "Freshness warning shown to user",     GREEN),
    ("U1", "Trust Gap",               "PM doesn't trust AI without sources",    "Every answer has source cards",       GREEN),
    ("U2", "Answer Length Mismatch",  "Too long for simple / too short for complex", "Intent-adaptive max_tokens (v2)", YELLOW),
    ("U3", "Follow-up Irrelevant",    "Suggestions don't match the query",     "Per-intent dynamic follow-ups",       GREEN),
]
y = Inches(1.8)
txbox(s, "ID", Inches(0.45), y, Inches(0.5), Inches(0.32), size=9, bold=True, color=DIM)
txbox(s, "Risk", Inches(1.0), y, Inches(2.0), Inches(0.32), size=9, bold=True, color=DIM)
txbox(s, "Description", Inches(3.1), y, Inches(3.5), Inches(0.32), size=9, bold=True, color=DIM)
txbox(s, "Mitigation", Inches(6.7), y, Inches(4.0), Inches(0.32), size=9, bold=True, color=DIM)
txbox(s, "Status", Inches(10.8), y, Inches(1.8), Inches(0.32), size=9, bold=True, color=DIM)
y += Inches(0.35)
hline(s, y)
y += Inches(0.04)

for rid, rname, desc, mitigation, col in risks:
    txbox(s, rid,        Inches(0.45), y, Inches(0.5), Inches(0.36), size=10, bold=True, color=col)
    txbox(s, rname,      Inches(1.0),  y, Inches(2.0), Inches(0.36), size=11, bold=True, color=WHITE)
    txbox(s, desc,       Inches(3.1),  y, Inches(3.5), Inches(0.36), size=10, color=DIM)
    txbox(s, mitigation, Inches(6.7),  y, Inches(4.0), Inches(0.36), size=10, color=DIM)
    status = "✓ Mitigated" if col == GREEN else "⚠ In Progress"
    txbox(s, status, Inches(10.8), y, Inches(1.8), Inches(0.36), size=10, bold=True, color=col)
    y += Inches(0.42)

txbox(s, "PM CHECKLIST  ·  Risk Management",
      Inches(0.45), Inches(6.6), Inches(12), Inches(0.32), size=10, bold=True, color=BLUE)
txbox(s, "✓  7 risks across AI, UX, and data layers   ✓  Each has owner + mitigation   ✓  Tracked as living doc",
      Inches(0.45), Inches(6.92), Inches(12), Inches(0.42), size=13, color=GREEN)

# ── S10: EVAL METHODOLOGY ────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 4 · EVAL METHODOLOGY", RGBColor(0xF8,0x71,0x71))
txbox(s, "How we proved it works — 3 scoring methods",
      Inches(0.55), Inches(0.85), Inches(12), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

methods = [
    ("A", "Automated\nLLM-as-Judge",
     "Claude Haiku scores every answer\non 3 dimensions against full source docs",
     "Fast · Consistent · Scales to 20 queries in minutes",
     BLUE),
    ("B", "Manual\nHuman Scoring",
     "PM scores same 20 queries independently\nGround truth baseline",
     "Catches what LLM misses · Adds business context",
     PURPLE),
    ("C", "Hybrid\nComparison",
     "Compare A vs B · Flag disagreements ≥2 pts\nas edge cases to investigate",
     "Identifies model blind spots · Informs v2 fixes",
     GREEN),
]
cx = Inches(0.55)
for letter, title, desc, why, col in methods:
    c = s.shapes.add_shape(1, cx, Inches(1.85), Inches(3.8), Inches(3.4))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    txbox(s, letter, cx + Inches(0.2), Inches(1.95), Inches(0.6), Inches(0.6), size=26, bold=True, color=col)
    txbox(s, title, cx + Inches(0.2), Inches(2.55), Inches(3.4), Inches(0.6), size=15, bold=True, color=WHITE)
    txbox(s, desc, cx + Inches(0.2), Inches(3.2), Inches(3.4), Inches(0.85), size=12, color=DIM)
    txbox(s, why, cx + Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.9), size=11, color=col)
    cx += Inches(4.2)

txbox(s, "20 queries · 5 categories · 3 metrics per query = 300 data points",
      Inches(0.55), Inches(5.45), Inches(12), Inches(0.45), size=15, bold=True, color=DIM, align=PP_ALIGN.CENTER)

txbox(s, "PM CHECKLIST  ·  Phase 4: Eval Design",
      Inches(0.55), Inches(6.1), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  Define eval before shipping   ✓  Mix automated + human scoring   ✓  Represent real user queries across 5 categories",
      Inches(0.55), Inches(6.45), Inches(12), Inches(0.45), size=13, color=GREEN)

# ── S11: EVAL RESULTS ────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 4 · RESULTS", RGBColor(0xF8,0x71,0x71))
txbox(s, "v1 Results — 5 of 6 metrics passed",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

result_cards = [
    ("Answer Relevance",  "4.8 / 5",  "Target ≥ 3.5",  GREEN),
    ("Citation Accuracy", "5.0 / 5",  "Target ≥ 4.0",  GREEN),
    ("Hallucination",     "4.9 / 5",  "Target ≥ 4.0",  GREEN),
    ("Retrieval Speed",   "1,397 ms", "Target < 4,000ms", GREEN),
    ("Total Latency",     "8,734 ms", "Target < 10,000ms", GREEN),
    ("Gen Latency",       "7,336 ms", "Target < 6,000ms",  RED),
]
cx = Inches(0.4)
row = 0
for i, (label, val, target, col) in enumerate(result_cards):
    if i == 3:
        cx = Inches(0.4)
        row = 1
    y_pos = Inches(1.75) if row == 0 else Inches(3.9)
    c = s.shapes.add_shape(1, cx, y_pos, Inches(3.9), Inches(1.8))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    txbox(s, label, cx + Inches(0.18), y_pos + Inches(0.15), Inches(3.5), Inches(0.38), size=11, bold=True, color=DIM)
    txbox(s, val,   cx + Inches(0.18), y_pos + Inches(0.52), Inches(3.5), Inches(0.55), size=28, bold=True, color=col)
    txbox(s, target, cx + Inches(0.18), y_pos + Inches(1.1), Inches(3.5), Inches(0.35), size=11, color=DIM)
    cx += Inches(4.3)

txbox(s, "✓ PASS — Seam v1 is production-quality on all quality metrics.\n⚠ Generation latency fails on complex multi-part queries (16–20s). Fixed in v2 with adaptive token capping.",
      Inches(0.55), Inches(6.15), Inches(12), Inches(0.75), size=13, color=WHITE)

# ── S12: EDGE CASES ──────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 4 · EDGE CASES & LEARNINGS", RGBColor(0xF8,0x71,0x17))
txbox(s, "4 edge cases — and what they taught us",
      Inches(0.55), Inches(0.85), Inches(12), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

edges = [
    ("Q2", "Permissions Acceptance Criteria",
     "LLM invented a Jira reference that isn't indexed",
     "Never speculate — only cite indexed sources", RED),
    ("Q3", "Walmart Escalation",
     "Slack data partial — LLM gave H=5, human gave H=3",
     "Flag Slack-heavy answers with completeness caveat", YELLOW),
    ("Q4", "Build vs Buy Analysis",
     "Ambiguous query — which product? LLM assumed correctly but shouldn't",
     "Trigger clarification chips for build vs buy queries", YELLOW),
    ("Q10", "Compare Review Engine vs Visual UGC",
     "No per-source attribution — hard to verify which claim came from where",
     "Prefix each subject's facts with [Source Title]:", YELLOW),
]
y = Inches(1.75)
for qid, qtitle, problem, fix, col in edges:
    c = s.shapes.add_shape(1, Inches(0.55), y, Inches(12.1), Inches(1.15))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    txbox(s, qid,    Inches(0.75), y + Inches(0.12), Inches(0.6), Inches(0.45), size=12, bold=True, color=col)
    txbox(s, qtitle, Inches(1.45), y + Inches(0.12), Inches(3.5), Inches(0.45), size=12, bold=True, color=WHITE)
    txbox(s, f"Problem: {problem}", Inches(5.1), y + Inches(0.1), Inches(4.1), Inches(0.45), size=11, color=DIM)
    txbox(s, f"Fix: {fix}", Inches(9.3), y + Inches(0.1), Inches(3.2), Inches(0.45), size=11, color=col)
    y += Inches(1.25)

txbox(s, "PM CHECKLIST  ·  Phase 4: Edge Cases",
      Inches(0.55), Inches(6.75), Inches(12), Inches(0.35), size=10, bold=True, color=BLUE)
txbox(s, "✓  Every edge case has a root cause   ✓  Maps to a specific risk from Phase 3   ✓  Each fix shipped in v2",
      Inches(0.55), Inches(7.1), Inches(12), Inches(0.35), size=13, color=GREEN)

# ── S13: V2 FIXES ────────────────────────────────────────────────────────────
s = blank_slide(prs)
phase_tag(s, "PHASE 4 · V2 IMPROVEMENTS", GREEN)
txbox(s, "What we fixed based on eval learnings",
      Inches(0.55), Inches(0.85), Inches(10), Inches(0.6),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.55))

fixes = [
    ("P0 · Critical", "'Not found' no longer guesses where info lives",
     "Eliminates hallucination risk on unanswered queries", GREEN),
    ("P1 · High",     "Build vs buy → ambiguity suggestions fired",
     "PM sees clarification chips, picks the right analysis", BLUE),
    ("P2 · High",     "Compare answers attribute each fact to its source",
     "[Review Engine PRD v3.0]: … / [Visual UGC Decision]: …", BLUE),
    ("P3 · Medium",   "Adaptive max_tokens: 500 (simple) / 800 (complex)",
     "~35% faster for single-intent queries", YELLOW),
    ("P4 · Medium",   "Slack-heavy answers flagged with completeness caveat",
     "PM knows when Slack data may be incomplete", YELLOW),
]
y = Inches(1.82)
for priority, fix, impact, col in fixes:
    c = s.shapes.add_shape(1, Inches(0.55), y, Inches(12.1), Inches(0.95))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    txbox(s, priority, Inches(0.75), y + Inches(0.1), Inches(2.0), Inches(0.4), size=11, bold=True, color=col)
    txbox(s, fix,      Inches(2.9),  y + Inches(0.1), Inches(5.1), Inches(0.4), size=12, color=WHITE)
    txbox(s, impact,   Inches(8.1),  y + Inches(0.1), Inches(4.4), Inches(0.4), size=11, color=DIM)
    y += Inches(1.04)

txbox(s, "PM CHECKLIST  ·  Iteration",
      Inches(0.55), Inches(7.1), Inches(12), Inches(0.32), size=10, bold=True, color=BLUE)
txbox(s, "✓  Every fix traces back to eval data   ✓  Prioritised by impact on user trust   ✓  All shipped and live in production",
      Inches(0.55), Inches(7.42), Inches(12), Inches(0.35), size=13, color=GREEN)

# ── S14: PM CHECKLIST MASTER ─────────────────────────────────────────────────
s = blank_slide(prs)
txbox(s, "The PM Checklist — All 4 Phases",
      Inches(0.55), Inches(0.3), Inches(12), Inches(0.65),
      size=28, bold=True, color=WHITE)
hline(s, Inches(1.02))

phases_checks = [
    ("Phase 1\nDiscovery", PURPLE, [
        "✓  Quantify the problem (numbers, not feelings)",
        "✓  Define primary persona with specificity",
        "✓  TAM / SAM / SOM with assumptions",
        "✓  Competitive analysis — name the gap",
    ]),
    ("Phase 2\nDesign", BLUE, [
        "✓  Explicit In / Out of scope list",
        "✓  Success metrics defined before building",
        "✓  Risk register with mitigation per risk",
        "✓  MVP delivers the core value loop end-to-end",
    ]),
    ("Phase 3\nBuild", GREEN, [
        "✓  Architecture decisions documented with why",
        "✓  Key trade-offs called out explicitly",
        "✓  Can explain tech in 60 seconds to a stakeholder",
        "✓  Risk register kept up to date",
    ]),
    ("Phase 4\nEval", RED, [
        "✓  Eval plan written before shipping",
        "✓  Mix automated + human + hybrid scoring",
        "✓  Every edge case has a root cause",
        "✓  v2 fixes traceable to eval data",
    ]),
]
cx = Inches(0.4)
for phase_label, col, checks_list in phases_checks:
    c = s.shapes.add_shape(1, cx, Inches(1.18), Inches(3.0), Inches(5.8))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    txbox(s, phase_label, cx + Inches(0.18), Inches(1.28), Inches(2.6), Inches(0.65), size=13, bold=True, color=col)
    cy = Inches(2.0)
    for chk in checks_list:
        txbox(s, chk, cx + Inches(0.18), cy, Inches(2.65), Inches(0.42), size=12, color=DIM)
        cy += Inches(0.5)
    cx += Inches(3.2)

txbox(s, "A great PM doesn't just ship — they show their work at every phase.",
      Inches(0.55), Inches(7.1), Inches(12), Inches(0.42), size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ── S15: LIVE DEMO SLIDE ─────────────────────────────────────────────────────
s = blank_slide(prs)
txbox(s, "Live Demo", Inches(1), Inches(2.5), Inches(11), Inches(1.2),
      size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "seam-ui-ebon.vercel.app",
      Inches(1), Inches(3.9), Inches(11), Inches(0.65),
      size=20, color=BLUE, align=PP_ALIGN.CENTER)

queries = [
    '"Why did Apex decide to build SAML 2.0 SSO?"',
    '"What did the build vs buy AI analysis conclude?"',
    '"Tell me about Apex — revenue, products, stakeholders"',
]
y = Inches(4.85)
for q in queries:
    txbox(s, f"→  {q}", Inches(2.5), y, Inches(8.3), Inches(0.42), size=14, color=DIM, italic=True)
    y += Inches(0.5)

# ── S16: CLOSING ─────────────────────────────────────────────────────────────
s = blank_slide(prs)
txbox(s, "seam.", Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
      size=72, bold=True, color=WHITE)
dot(s, Inches(4.47), Inches(2.12), r=14)
txbox(s, "One query. Every answer. All your tools.",
      Inches(1.5), Inches(3.15), Inches(9), Inches(0.65),
      size=22, color=DIM)
hline(s, Inches(3.9))
txbox(s, "Built by: Gunjan  ·  Capstone Cohort 6  ·  Deployed at seam-ui-ebon.vercel.app",
      Inches(1.5), Inches(4.1), Inches(10), Inches(0.5), size=14, color=DIM)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "Seam_Capstone_Deck.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
